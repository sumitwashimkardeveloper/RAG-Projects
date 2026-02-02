import os
import json
from abc import ABC, abstractmethod
from typing import List, Dict, Any, Optional
from pathlib import Path
import base64
from io import BytesIO

from PIL import Image
import cv2
import numpy as np

try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

try:
    import librosa
    LIBROSA_AVAILABLE = True
except ImportError:
    LIBROSA_AVAILABLE = False

try:
    from moviepy.editor import VideoFileClip
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

from config import Config

class MediaProcessor(ABC):
    def __init__(self):
        self.config = Config()
        os.makedirs(self.config.TEMP_DIR, exist_ok=True)

    @abstractmethod
    def process(self, file_path: str) -> Dict[str, Any]:
        pass

    @abstractmethod
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        pass

    def chunk_content(self, content: str, chunk_size: int = None) -> List[str]:
        if chunk_size is None:
            chunk_size = self.config.CHUNK_SIZE

        chunks = []
        for i in range(0, len(content), chunk_size - self.config.CHUNK_OVERLAP):
            chunk = content[i:i + chunk_size]
            if chunk.strip():
                chunks.append(chunk)
        return chunks


class ImageProcessor(MediaProcessor):
    def process(self, file_path: str) -> Dict[str, Any]:
        metadata = self.extract_metadata(file_path)

        image = Image.open(file_path)
        image_array = np.array(image)

        text_content = ""
        if self.config.OCR_ENABLED and TESSERACT_AVAILABLE:
            text_content = pytesseract.image_to_string(image)

        objects = {}
        if self.config.OBJECT_DETECTION_ENABLED:
            objects = self._detect_objects(image_array)

        image_base64 = self._encode_image(file_path)

        return {
            "type": "image",
            "file_path": file_path,
            "metadata": metadata,
            "text_content": text_content,
            "objects": objects,
            "image_base64": image_base64,
            "chunks": self.chunk_content(text_content) if text_content else []
        }

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        image = Image.open(file_path)
        return {
            "format": image.format,
            "size": image.size,
            "mode": image.mode,
            "file_size_bytes": os.path.getsize(file_path),
            "dpi": image.info.get("dpi", None)
        }

    def _detect_objects(self, image_array: np.ndarray) -> Dict[str, Any]:
        gray = cv2.cvtColor(image_array, cv2.COLOR_RGB2GRAY)
        contours, _ = cv2.findContours(gray, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)

        objects = {
            "detected_count": len(contours),
            "contours_info": []
        }

        for contour in contours[:10]:
            area = cv2.contourArea(contour)
            if area > 100:
                x, y, w, h = cv2.boundingRect(contour)
                objects["contours_info"].append({
                    "area": float(area),
                    "bbox": {"x": int(x), "y": int(y), "w": int(w), "h": int(h)}
                })

        return objects

    def _encode_image(self, file_path: str) -> str:
        with open(file_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')


class TableProcessor(MediaProcessor):
    def process(self, file_path: str) -> Dict[str, Any]:
        metadata = self.extract_metadata(file_path)

        if file_path.endswith('.csv'):
            import pandas as pd
            df = pd.read_csv(file_path)
            table_data = df.to_dict('records')
            text_content = df.to_string()
        else:
            table_data = []
            text_content = ""

        return {
            "type": "table",
            "file_path": file_path,
            "metadata": metadata,
            "table_data": table_data,
            "text_content": text_content,
            "chunks": self.chunk_content(text_content)
        }

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        import pandas as pd
        df = pd.read_csv(file_path) if file_path.endswith('.csv') else pd.DataFrame()
        return {
            "rows": len(df),
            "columns": len(df.columns),
            "column_names": list(df.columns),
            "file_size_bytes": os.path.getsize(file_path)
        }


class ChartProcessor(MediaProcessor):
    def process(self, file_path: str) -> Dict[str, Any]:
        metadata = self.extract_metadata(file_path)

        image = Image.open(file_path)
        image_array = np.array(image)

        text_content = ""
        if TESSERACT_AVAILABLE:
            text_content = pytesseract.image_to_string(image)

        chart_data = self._extract_chart_elements(image_array)

        return {
            "type": "chart",
            "file_path": file_path,
            "metadata": metadata,
            "text_content": text_content,
            "chart_elements": chart_data,
            "chunks": self.chunk_content(text_content + " " + json.dumps(chart_data))
        }

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        image = Image.open(file_path)
        return {
            "format": image.format,
            "size": image.size,
            "file_size_bytes": os.path.getsize(file_path)
        }

    def _extract_chart_elements(self, image_array: np.ndarray) -> Dict[str, Any]:
        hsv = cv2.cvtColor(image_array, cv2.COLOR_RGB2HSV)

        colors_found = []
        for lower, upper in [
            (np.array([0, 100, 100]), np.array([10, 255, 255])),
            (np.array([110, 100, 100]), np.array([130, 255, 255]))
        ]:
            mask = cv2.inRange(hsv, lower, upper)
            if cv2.countNonZero(mask) > 0:
                colors_found.append("color_detected")

        return {
            "chart_type": "unknown",
            "colors_detected": len(colors_found),
            "dimensions": image_array.shape[:2]
        }


class VideoProcessor(MediaProcessor):
    def process(self, file_path: str) -> Dict[str, Any]:
        metadata = self.extract_metadata(file_path)

        frames_data = []
        transcript = ""

        if MOVIEPY_AVAILABLE:
            video = VideoFileClip(file_path)

            frame_step = int(video.fps / self.config.VIDEO_FRAME_SAMPLING)
            frame_count = 0

            for frame in video.iter_frames(fps=self.config.VIDEO_FRAME_SAMPLING):
                frame_count += 1
                frame_description = self._describe_frame(frame)
                frames_data.append({
                    "frame_number": frame_count,
                    "description": frame_description
                })

            if metadata.get("has_audio"):
                transcript = self._extract_audio_transcript(file_path)

            video.close()

        all_text = " ".join([f["description"] for f in frames_data]) + " " + transcript

        return {
            "type": "video",
            "file_path": file_path,
            "metadata": metadata,
            "frames": frames_data,
            "transcript": transcript,
            "text_content": all_text,
            "chunks": self.chunk_content(all_text)
        }

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        if not MOVIEPY_AVAILABLE:
            return {"error": "moviepy not available"}

        video = VideoFileClip(file_path)
        metadata = {
            "duration_seconds": float(video.duration),
            "fps": float(video.fps),
            "resolution": video.size,
            "file_size_bytes": os.path.getsize(file_path),
            "has_audio": video.audio is not None
        }
        video.close()
        return metadata

    def _describe_frame(self, frame: np.ndarray) -> str:
        gray = cv2.cvtColor(frame, cv2.COLOR_RGB2GRAY)
        edge_count = cv2.Canny(gray, 100, 200).sum()

        return f"Frame with edge intensity: {edge_count}"

    def _extract_audio_transcript(self, file_path: str) -> str:
        if LIBROSA_AVAILABLE:
            y, sr = librosa.load(file_path)
            return f"Audio loaded with sample rate {sr}"
        return ""


class AudioProcessor(MediaProcessor):
    def process(self, file_path: str) -> Dict[str, Any]:
        metadata = self.extract_metadata(file_path)

        transcript = ""
        audio_features = {}

        if LIBROSA_AVAILABLE:
            y, sr = librosa.load(file_path)
            audio_features = self._extract_audio_features(y, sr)
            transcript = self._transcribe_audio(file_path)

        text_content = transcript + " " + json.dumps(audio_features)

        return {
            "type": "audio",
            "file_path": file_path,
            "metadata": metadata,
            "transcript": transcript,
            "audio_features": audio_features,
            "text_content": text_content,
            "chunks": self.chunk_content(text_content)
        }

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        if not LIBROSA_AVAILABLE:
            return {"file_size_bytes": os.path.getsize(file_path)}

        y, sr = librosa.load(file_path)
        return {
            "duration_seconds": float(librosa.get_duration(y=y, sr=sr)),
            "sample_rate": int(sr),
            "channels": 1 if len(y.shape) == 1 else y.shape[0],
            "file_size_bytes": os.path.getsize(file_path)
        }

    def _extract_audio_features(self, y: np.ndarray, sr: int) -> Dict[str, Any]:
        return {
            "rms_energy": float(np.sqrt(np.mean(y**2))),
            "zero_crossing_rate": float(np.mean(librosa.feature.zero_crossing_rate(y))),
            "spectral_centroid": float(np.mean(librosa.feature.spectral_centroid(y=y, sr=sr))),
            "mfcc_mean": float(np.mean(librosa.feature.mfcc(y=y, sr=sr, n_mfcc=13)))
        }

    def _transcribe_audio(self, file_path: str) -> str:
        return f"Audio transcription would happen here for {Path(file_path).name}"


class PowerPointProcessor(MediaProcessor):
    def process(self, file_path: str) -> Dict[str, Any]:
        metadata = self.extract_metadata(file_path)

        slides_data = []
        all_text = ""

        if PPTX_AVAILABLE:
            prs = Presentation(file_path)

            for slide_idx, slide in enumerate(prs.slides):
                slide_content = {
                    "slide_number": slide_idx + 1,
                    "text": "",
                    "shapes": []
                }

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content["text"] += shape.text + " "

                    if shape.has_table:
                        table_data = self._extract_table_from_shape(shape)
                        slide_content["shapes"].append({
                            "type": "table",
                            "data": table_data
                        })

                if hasattr(slide, "notes_slide"):
                    notes = slide.notes_slide.notes_text_frame.text
                    slide_content["notes"] = notes

                slides_data.append(slide_content)
                all_text += slide_content["text"] + " "

        return {
            "type": "powerpoint",
            "file_path": file_path,
            "metadata": metadata,
            "slides": slides_data,
            "text_content": all_text,
            "chunks": self.chunk_content(all_text)
        }

    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        if not PPTX_AVAILABLE:
            return {"file_size_bytes": os.path.getsize(file_path)}

        prs = Presentation(file_path)
        return {
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "file_size_bytes": os.path.getsize(file_path)
        }

    def _extract_table_from_shape(self, shape) -> List[Dict[str, Any]]:
        if not shape.has_table:
            return []

        table = shape.table
        table_data = []

        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            table_data.append(row_data)

        return table_data


class MediaProcessorFactory:
    @staticmethod
    def get_processor(file_path: str) -> Optional[MediaProcessor]:
        ext = Path(file_path).suffix.lower()

        if ext in Config.SUPPORTED_IMAGE_FORMATS:
            return ImageProcessor()
        elif ext in {".csv", ".xlsx"}:
            return TableProcessor()
        elif ext in Config.SUPPORTED_VIDEO_FORMATS:
            return VideoProcessor()
        elif ext in Config.SUPPORTED_AUDIO_FORMATS:
            return AudioProcessor()
        elif ext == ".pptx":
            return PowerPointProcessor()
        elif ext == ".pdf":
            return ImageProcessor()

        return None
