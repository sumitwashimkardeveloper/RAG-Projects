"""
Document Processors for extracting text from various file formats
"""

from abc import ABC, abstractmethod
from typing import Optional
import logging

logger = logging.getLogger(__name__)


class BaseDocumentProcessor(ABC):
    """Base class for document processors"""

    @abstractmethod
    def extract_text(self, file_path: str) -> Optional[str]:
        """
        Extract text from document

        Args:
            file_path: Path to document file

        Returns:
            Extracted text or None if extraction fails
        """
        pass


class PDFProcessor(BaseDocumentProcessor):
    """Extract text from PDF files"""

    def extract_text(self, file_path: str) -> Optional[str]:
        """Extract text from PDF"""
        try:
            import PyPDF2

            text_parts = []
            with open(file_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(reader.pages):
                    try:
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)
                    except Exception as e:
                        logger.warning(f"Failed to extract page {page_num} from {file_path}: {e}")

            return "\n".join(text_parts) if text_parts else None
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
            return None


class DOCXProcessor(BaseDocumentProcessor):
    """Extract text from DOCX files"""

    def extract_text(self, file_path: str) -> Optional[str]:
        """Extract text from DOCX"""
        try:
            from docx import Document

            doc = Document(file_path)
            text_parts = []

            # Extract from paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        [cell.text.strip() for cell in row.cells]
                    )
                    if row_text.strip():
                        text_parts.append(row_text)

            return "\n".join(text_parts) if text_parts else None
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {str(e)}")
            return None


class ExcelProcessor(BaseDocumentProcessor):
    """Extract text from Excel files"""

    def extract_text(self, file_path: str) -> Optional[str]:
        """Extract text from Excel"""
        try:
            import openpyxl

            text_parts = []
            wb = openpyxl.load_workbook(file_path, data_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")

                for row in sheet.iter_rows(values_only=True):
                    row_values = [
                        str(cell) if cell is not None else ""
                        for cell in row
                    ]
                    row_text = " | ".join(row_values).strip()
                    if row_text and row_text != " | " * (len(row) - 1):
                        text_parts.append(row_text)

            return "\n".join(text_parts) if text_parts else None
        except Exception as e:
            logger.error(f"Error processing Excel {file_path}: {str(e)}")
            return None


class CSVProcessor(BaseDocumentProcessor):
    """Extract text from CSV files"""

    def extract_text(self, file_path: str) -> Optional[str]:
        """Extract text from CSV"""
        try:
            import pandas as pd

            df = pd.read_csv(file_path)
            # Convert dataframe to string with formatting
            return df.to_string()
        except Exception as e:
            logger.error(f"Error processing CSV {file_path}: {str(e)}")
            return None


class MarkdownProcessor(BaseDocumentProcessor):
    """Extract text from Markdown files"""

    def extract_text(self, file_path: str) -> Optional[str]:
        """Extract text from Markdown"""
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error processing Markdown {file_path}: {str(e)}")
            return None


class PlainTextProcessor(BaseDocumentProcessor):
    """Extract text from plain text files"""

    def extract_text(self, file_path: str) -> Optional[str]:
        """Extract text from plain text file"""
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {str(e)}")
            return None


class PowerPointProcessor(BaseDocumentProcessor):
    """Extract text from PowerPoint files"""

    def extract_text(self, file_path: str) -> Optional[str]:
        """Extract text from PowerPoint"""
        try:
            from pptx import Presentation

            text_parts = []
            presentation = Presentation(file_path)

            for slide_num, slide in enumerate(presentation.slides, 1):
                text_parts.append(f"Slide {slide_num}:")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

            return "\n".join(text_parts) if text_parts else None
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {str(e)}")
            return None
